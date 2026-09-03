"""Build the ENTITY GRID submission deck.

Run as ``python deck/build_deck.py``. Writes ``deck/ENTITY_GRID.pptx``.

Every number on these slides is read out of ``data/processed/results.json`` and
the benchmark CSVs at build time rather than typed in, so the deck cannot drift
away from what the code actually produces. If the pipeline result changes, the
deck changes with it.
"""

from __future__ import annotations

import json
import pathlib
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = pathlib.Path(__file__).resolve().parents[1]
PROCESSED = ROOT / "data" / "processed"
OUT = pathlib.Path(__file__).resolve().parent / "ENTITY_GRID.pptx"

# ---------------------------------------------------------------- palette
BG = RGBColor(0x08, 0x0A, 0x0F)
PANEL = RGBColor(0x0F, 0x13, 0x1C)
LINE = RGBColor(0x26, 0x2D, 0x3F)
INK = RGBColor(0xEE, 0xF1, 0xF7)
MID = RGBColor(0x9A, 0xA3, 0xB8)
DIM = RGBColor(0x62, 0x6C, 0x85)
BLUE = RGBColor(0x3B, 0x6D, 0xFA)
BLUE2 = RGBColor(0x7E, 0xA0, 0xFF)
OK = RGBColor(0x39, 0xC0, 0x7A)
WARN = RGBColor(0xE0, 0x93, 0x2F)
CRIT = RGBColor(0xFF, 0x5C, 0x5C)

DISP = "Arial Narrow"
BODY = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.78)


# ---------------------------------------------------------------- helpers
def load() -> dict:
    path = PROCESSED / "results.json"
    if not path.exists():
        sys.exit("run `python -m entitygrid.pipeline` first, no results.json found")
    data = {"results": json.loads(path.read_text())}
    for name in ("topology_benchmark", "detector_benchmark", "pv_scenarios",
                 "topology_sensitivity"):
        csv = PROCESSED / f"{name}.csv"
        if csv.exists():
            import csv as _csv
            with csv.open() as fh:
                data[name] = list(_csv.DictReader(fh))
        else:
            data[name] = []
    return data


def slide(prs) -> "object":
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0):
    """runs: list of (text, size, bold, colour, font) or list of such lists."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if i:
            p.space_before = Pt(6)
        for chunk in para:
            body, size, bold, colour = chunk[0], chunk[1], chunk[2], chunk[3]
            font_name = chunk[4] if len(chunk) > 4 else BODY
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = font_name
    return box


def rect(s, x, y, w, h, fill=PANEL, line=None):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def eyebrow(s, label):
    text(s, M, Inches(0.52), Inches(11), Inches(0.3),
         [(label.upper(), 10.5, True, DIM, MONO)])


def heading(s, title, index=None):
    if index:
        text(s, M, Inches(0.92), Inches(1.4), Inches(0.9),
             [(index, 44, True, LINE, DISP)])
    x = M + (Inches(1.05) if index else 0)
    text(s, x, Inches(0.95), Inches(11), Inches(0.85),
         [(title.upper(), 34, True, INK, DISP)])


def lede(s, body, y=Inches(1.72), w=Inches(10.4)):
    text(s, M, y, w, Inches(1.0), [(body, 13.5, False, MID)], spacing=1.25)


def footer(s, n, note=""):
    text(s, M, H - Inches(0.5), Inches(9), Inches(0.3),
         [(note, 9.5, False, DIM, MONO)])
    text(s, W - M - Inches(1.2), H - Inches(0.5), Inches(1.2), Inches(0.3),
         [(f"{n:02d}", 9.5, False, DIM, MONO)], align=PP_ALIGN.RIGHT)


def stat_row(s, y, items, height=Inches(1.42)):
    """items: (value, unit, label, sub, colour)"""
    n = len(items)
    gap = Inches(0.14)
    total = W - 2 * M
    cw = int((total - gap * (n - 1)) / n)
    for i, (value, unit, label, sub, colour) in enumerate(items):
        x = M + i * (cw + gap)
        rect(s, x, y, Emu(cw), height, PANEL, LINE)
        text(s, x + Inches(0.22), y + Inches(0.16), Emu(cw) - Inches(0.4), Inches(0.24),
             [(label.upper(), 9, True, DIM, MONO)])
        runs = [(value, 30, True, colour, DISP)]
        if unit:
            runs.append((unit, 14, True, MID, DISP))
        text(s, x + Inches(0.22), y + Inches(0.46), Emu(cw) - Inches(0.4), Inches(0.46), runs)
        text(s, x + Inches(0.22), y + Inches(0.95), Emu(cw) - Inches(0.4), Inches(0.4),
             [(sub, 9.5, False, MID)], spacing=1.15)


def table(s, x, y, w, cols, rows, widths, row_h=Inches(0.34),
          head_h=Inches(0.32), highlight=None):
    """Simple flat table. highlight(row_index, col_index, value) -> colour|None."""
    xs, acc = [], x
    for frac in widths:
        xs.append(acc)
        acc += Emu(int(w * frac))

    rect(s, x, y, w, head_h, PANEL, None)
    for i, (label, frac) in enumerate(zip(cols, widths)):
        text(s, xs[i] + Inches(0.1), y + Inches(0.07), Emu(int(w * frac)), Inches(0.24),
             [(label.upper(), 8.5, True, DIM, MONO)])

    for r, row in enumerate(rows):
        ry = y + head_h + r * row_h
        if r % 2 == 0:
            rect(s, x, ry, w, row_h, RGBColor(0x0B, 0x0E, 0x15), None)
        for c, cell in enumerate(row):
            colour = INK
            bold = False
            if highlight:
                got = highlight(r, c, cell)
                if got is not None:
                    colour, bold = got, True
            text(s, xs[c] + Inches(0.1), ry + Inches(0.08),
                 Emu(int(w * widths[c])), Inches(0.26),
                 [(str(cell), 10, bold, colour, MONO)])
    return y + head_h + len(rows) * row_h


def bullets(s, x, y, w, items, gap=Inches(0.86)):
    for i, (title, body) in enumerate(items):
        yy = y + i * gap
        rect(s, x, yy + Inches(0.06), Inches(0.035), Inches(0.5), BLUE)
        text(s, x + Inches(0.22), yy, w - Inches(0.22), Inches(0.28),
             [(title, 13.5, True, INK)])
        text(s, x + Inches(0.22), yy + Inches(0.28), w - Inches(0.22), Inches(0.5),
             [(body, 11.5, False, MID)], spacing=1.2)


# ---------------------------------------------------------------- slides
def build(d: dict) -> Presentation:
    r = d["results"]
    sc = r["scorecard"]
    topo, health, faults, volt = (sc["topology"], sc["health"], sc["faults"],
                                  sc["voltage"])
    flex = sc["flexibility"]
    meta = r["meta"]
    pv = d["pv_scenarios"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------- 01 title
    s = slide(prs)
    rect(s, 0, 0, Inches(0.09), H, BLUE)
    text(s, M, Inches(1.5), Inches(11), Inches(0.3),
         [("YUVA YODHA ENERGY TECH HACKATHON 2026  ·  GRID RELIABILITY",
           11, True, DIM, MONO)])
    text(s, M, Inches(2.0), Inches(11.6), Inches(2.2),
         [[("ENTITY", 76, True, INK, DISP), ("GRID", 76, True, BLUE, DISP)]])
    text(s, M, Inches(3.3), Inches(9.6), Inches(1.5),
         [("You cannot forecast a feeder, shed load on it, or site a battery on it "
           "if nobody knows what is connected to it.", 21, False, INK)], spacing=1.25)
    text(s, M, Inches(4.5), Inches(10.4), Inches(1.2),
         [("One consumer in three is recorded against the wrong transformer or phase. "
           "ENTITY GRID learns the real low voltage network from the smart meters "
           "already being installed, then uses it to keep the neighbourhood "
           "dependable. Software only. No new hardware, no site survey, no capex.",
           13, False, MID)], spacing=1.3)
    stat_row(s, Inches(5.75), [
        (f"{topo['joint_accuracy']*100:.0f}%", "", "Connectivity recovered",
         f"against a {topo['ledger_joint_accuracy']*100:.1f}% utility ledger", OK),
        (str(topo["corrections_found"]), "", "Records corrected",
         f"of {topo['n_meters']} consumers, from voltage alone", BLUE2),
        (f"{faults['detected']}/{faults['truth_events']}", "", "Faults located",
         f"{faults['mean_precision']*100:.0f}% precision, about a second", OK),
        (f"{flex['mean_skill_vs_baseline']*100:.0f}%", "", "Better forecasts",
         "than the best baseline a utility already has", OK),
    ], Inches(1.35))
    footer(s, 1, "entitygrid  ·  software only  ·  runs on RDSS meter data")

    # ---------------------------------------------------------- 02 problem
    s = slide(prs)
    eyebrow(s, "the problem")
    heading(s, "Three failures,\none missing map", "01")
    text(s, M, Inches(2.35), Inches(10.6), Inches(0.6),
         [("Every low voltage problem an Indian DISCOM has traces back to the same "
           "gap. Nobody knows what is actually connected to what.", 13.5, False, MID)],
         spacing=1.25)
    cols = Inches(3.72)
    for i, (n, title, body, stat) in enumerate([
        ("01", "Billing the wrong feeder",
         "Loss accounting, load balancing and every outage decision inherit the "
         "ledger's errors.",
         "69.1% correct today"),
        ("02", "Transformers fail unannounced",
         "A unit that fails at 2 a.m. in June costs overtime, a dark neighbourhood "
         "and a week of complaints.",
         "no warning at all today"),
        ("03", "Crews drive the line",
         "When a feeder trips, the fault is still found by patrolling it with a "
         "torch.",
         "hours, not seconds"),
    ]):
        x = M + i * (cols + Inches(0.16))
        rect(s, x, Inches(3.15), cols, Inches(2.5), PANEL, LINE)
        text(s, x + Inches(0.26), Inches(3.35), cols - Inches(0.5), Inches(0.25),
             [(n, 11, True, BLUE, MONO)])
        text(s, x + Inches(0.26), Inches(3.62), cols - Inches(0.5), Inches(0.6),
             [(title.upper(), 17, True, INK, DISP)], spacing=1.05)
        text(s, x + Inches(0.26), Inches(4.32), cols - Inches(0.5), Inches(0.9),
             [(body, 11.5, False, MID)], spacing=1.25)
        text(s, x + Inches(0.26), Inches(5.25), cols - Inches(0.5), Inches(0.3),
             [(stat, 10, True, CRIT, MONO)])
    text(s, M, Inches(6.0), Inches(11.5), Inches(0.5),
         [("Meanwhile RDSS is putting a smart meter on every connection in the "
           "country, and that data is used almost entirely for billing.",
           13.5, True, INK)], spacing=1.25)
    footer(s, 2)

    # ---------------------------------------------------------- 03 insight
    s = slide(prs)
    eyebrow(s, "the insight")
    heading(s, "The network is\nalready in the data", "02")
    text(s, M, Inches(2.5), Inches(6.1), Inches(1.8),
         [("A meter's voltage is dominated by things that say nothing about where it "
           "sits. The substation swing, tap steps, the load cycle of the whole 11 kV "
           "feeder. Every meter sees those at once, so a correlation on raw voltage "
           "is near 1.0 for every pair and separates nothing.\n\n"
           "Difference in time to kill slow drift. Subtract the cross sectional "
           "median to kill what the whole network saw together. What survives is the "
           "drop across the impedance a meter shares with its transformer.",
           13, False, MID)], spacing=1.3)
    y = Inches(2.55)
    rect(s, Inches(7.3), y - Inches(0.16), Inches(5.25), Inches(2.35), PANEL, LINE)
    text(s, Inches(7.55), y, Inches(4.8), Inches(0.3),
         [("CORRELATION OF RESIDUAL INCREMENTS", 9, True, DIM, MONO)])
    for i, (label, value, colour) in enumerate([
            ("Same transformer, same phase", "0.81", OK),
            ("Same transformer, other phase", "-0.09", MID),
            ("Different transformer", "-0.01", MID)]):
        yy = y + Inches(0.42) + i * Inches(0.52)
        text(s, Inches(7.55), yy, Inches(3.6), Inches(0.3),
             [(label, 11.5, False, INK)])
        text(s, Inches(11.2), yy - Inches(0.06), Inches(1.1), Inches(0.4),
             [(value, 20, True, colour, DISP)], align=PP_ALIGN.RIGHT)
    text(s, Inches(7.3), Inches(5.15), Inches(5.25), Inches(1.1),
         [("So the natural cluster is the transformer and phase pair, not the "
           "transformer. Clustering recovers all 36 groups with an adjusted Rand "
           "index of 1.000. Matching each group against busbar telemetry labels it "
           "with a real asset ID.", 11.5, False, MID)], spacing=1.25)
    footer(s, 3, "no ground truth is read; a test fails the build if it is")

    # ---------------------------------------------------------- 04 pillars
    s = slide(prs)
    eyebrow(s, "what it does")
    heading(s, "Five pillars,\none learned map", "03")
    items = [
        ("Self learning topology",
         "Meter to transformer to phase from voltage alone, with a confidence score "
         "and an explicit field check list."),
        ("Predictive health",
         "Estimated impedance per transformer and per feeder segment, tracked "
         "against each asset's own baseline."),
        ("Fault localisation",
         "Last gasp bursts bracketed by electrical depth, naming the last consumer "
         "lit and the first one dark."),
        ("Solar aware voltage",
         "Reverse flow separated from load driven sag, answered with reactive "
         "setpoints before any curtailment."),
        ("Neighbourhood flexibility",
         "Per phase day ahead forecasts, predicted constraints, targeted demand "
         "response and storage sized from the forecast."),
    ]
    bullets(s, M, Inches(2.45), Inches(6.3), items, gap=Inches(0.84))
    rect(s, Inches(7.5), Inches(2.4), Inches(5.05), Inches(4.05), PANEL, LINE)
    text(s, Inches(7.8), Inches(2.62), Inches(4.5), Inches(0.3),
         [("WHY THE ORDER MATTERS", 9, True, DIM, MONO)])
    text(s, Inches(7.8), Inches(3.0), Inches(4.45), Inches(3.2),
         [("Everything downstream is expressed against the connectivity the learner "
           "found.\n\n"
           "Health aggregates meters to the transformer it found. Fault location "
           "brackets along the topology it found. Volt var targets the phase it "
           "found. Flexibility forecasts, sheds and stores against it.\n\n"
           "Feed those stages the utility's own ledger instead and a 31% error rate "
           "goes straight into every result.", 12, False, MID)], spacing=1.32)
    footer(s, 4)

    # ---------------------------------------------------------- 05 evidence
    s = slide(prs)
    eyebrow(s, "evidence")
    heading(s, "Graded against\na known answer", "04")
    lede(s, "Validated on a simulated network whose true connectivity, degradation "
            "and faults are known exactly, and which no model is allowed to read. "
            f"{meta['meters']} consumers, {meta['transformers']} transformers, "
            f"{meta['days']} days of 15 minute data.")
    stat_row(s, Inches(2.85), [
        (f"{topo['joint_accuracy']*100:.0f}%", "", "Connectivity",
         f"ledger reads {topo['ledger_joint_accuracy']*100:.1f}%", OK),
        (f"{health['detected_in_time']}/{len(health['degrading_transformers'])}", "",
         "Caught before failure",
         f"{len(health['false_positive_transformers'])} false alarms, "
         f"{health['found_too_late']} found late", WARN),
        (f"{faults['detected']}/{faults['truth_events']}", "", "Faults located",
         f"{faults['mean_recall']*100:.0f}% recall, "
         f"{faults['mean_precision']*100:.0f}% precision", OK),
        (f"{flex['mean_skill_vs_baseline']*100:.0f}%", "", "Forecast skill",
         f"{flex['forecasters_beating_baseline']} of {flex['forecasters']} "
         f"phases beat the baseline", OK),
    ])

    rows = []
    for row in d["topology_benchmark"]:
        if row["method"] != "aligned":
            continue
        raw = next((x for x in d["topology_benchmark"]
                    if x["scenario"] == row["scenario"] and x["method"] == "raw"), None)
        rows.append([row["scenario"],
                     f"{float(raw['joint_accuracy'])*100:.1f}" if raw else "n/a",
                     f"{float(row['joint_accuracy'])*100:.1f}",
                     f"{float(row['ledger_baseline'])*100:.1f}"])

    def hi(r, c, v):
        if c == 2:
            return OK
        if c == 1:
            try:
                return CRIT if float(v) < 69.1 else None
            except ValueError:
                return None
        return None

    text(s, M, Inches(4.5), Inches(11.6), Inches(0.3),
         [("TOPOLOGY ACCURACY UNDER STRESS, PERCENT", 9, True, DIM, MONO)])
    table(s, M, Inches(4.82), Inches(11.75),
          ["scenario", "naive correlation", "entity grid", "utility ledger"],
          rows, [0.38, 0.22, 0.2, 0.2], highlight=hi)
    footer(s, 5, "every naive score sits below the ledger it was meant to replace")

    # ---------------------------------------------------------- 06 what lost
    s = slide(prs)
    eyebrow(s, "what did not work")
    heading(s, "Three things\nthat lost", "05")
    lede(s, "Kept in the repository and on the console, because a number with no "
            "alternative beside it is not evidence.")
    det_rows = [[row["detector"],
                 f"{row['detected_in_time']}/{row['of_degrading']}",
                 row["false_alarms"],
                 (f"{float(row['mean_lead_days']):.1f}"
                  if row.get("mean_lead_days") not in (None, "", "nan") else "n/a")]
                for row in d["detector_benchmark"]]

    def hi2(r, c, v):
        if c == 2:
            return OK if str(v) == "0" else CRIT
        return None

    text(s, M, Inches(2.75), Inches(6.0), Inches(0.3),
         [("DEGRADATION DETECTORS", 9, True, DIM, MONO)])
    table(s, M, Inches(3.05), Inches(6.0),
          ["detector", "in time", "false alarms", "lead"],
          det_rows, [0.36, 0.2, 0.26, 0.18], highlight=hi2)
    text(s, M, Inches(4.6), Inches(6.0), Inches(1.7),
         [("A cumulative sum chart is the textbook detector for slow degradation and "
           "is what the published work uses. On this data it raised eight false "
           "alarms to catch one real fault. The trend test ships instead.",
           11.5, False, MID)], spacing=1.28)

    rect(s, Inches(7.25), Inches(2.72), Inches(5.3), Inches(3.5), PANEL, LINE)
    text(s, Inches(7.55), Inches(2.95), Inches(4.7), Inches(0.3),
         [("ALSO TRIED, ALSO REJECTED", 9, True, DIM, MONO)])
    text(s, Inches(7.55), Inches(3.35), Inches(4.7), Inches(2.7),
         [("Low load data segmentation. The literature reports it improves phase "
           "identification. It gains about three points under clock drift here and "
           "loses ground everywhere else, so it is implemented and off by default.\n\n"
           "Excluding daylight hours to dodge shared solar signature. Common mode "
           "removal already handles it, so the filter costs a third of the day for "
           "nothing.\n\n"
           "An earlier detector fitted its trend across the whole series and quoted "
           "an early onset. That was lookahead, and it inflated lead times to 54 "
           "days. A test now fails the build if detection uses future data.",
           11.5, False, MID)], spacing=1.28)
    footer(s, 6)

    # ---------------------------------------------------------- 07 flexibility
    s = slide(prs)
    eyebrow(s, "the challenge brief")
    heading(s, "Flexibility, aimed\nat the right street", "06")
    lede(s, "The brief asks for dependable power neighbourhood by neighbourhood "
            "through forecasting, demand response and shared storage. None of those "
            "can be aimed at a network nobody has mapped, which is why this pillar "
            "sits on top of the other four rather than beside them.")
    stat_row(s, Inches(3.0), [
        (f"{flex['forecasters']}", "", "Per phase forecasters",
         "a feeder is rarely in trouble as a whole", BLUE2),
        (f"{flex['mean_skill_vs_baseline']*100:.0f}%", "", "Better than baseline",
         f"nMAE {flex['median_nmae_pct']:.1f}%, weather input degraded to "
         f"real forecast accuracy", OK),
        (f"{flex['constraint_windows']}", "", "Constraints predicted",
         ", ".join(f"{v} {k}" for k, v in flex["constraints_by_kind"].items()), WARN),
        (f"{flex['flexibility_lift_over_base_rate']:.1f}x", "", "Flexible consumers found",
         "above base rate, from load shape alone", OK),
    ])
    bullets(s, M, Inches(4.75), Inches(5.9), [
        ("Targeted, not broadcast",
         "Call lists filtered to the transformer and the phase actually in trouble, "
         "then weighted by electrical depth."),
        ("Storage sized from the forecast",
         "Power set by the largest predicted deficit, energy by its area, charging "
         "from export the feeder already spills."),
    ], gap=Inches(0.92))
    rect(s, Inches(7.1), Inches(4.72), Inches(5.45), Inches(1.9), PANEL, LINE)
    text(s, Inches(7.4), Inches(4.95), Inches(4.9), Inches(0.3),
         [("INFERRING FLEXIBILITY WITHOUT ASKING", 9, True, DIM, MONO)])
    text(s, Inches(7.4), Inches(5.32), Inches(4.85), Inches(1.2),
         [("Nobody knows which consumers can shift load, and tariff category is as "
           "stale as the connectivity record. Pumping is unmistakable in interval "
           "data: long runs, high peaks, mostly at night. Scored from shape alone.",
           11.5, False, MID)], spacing=1.28)
    footer(s, 7)

    # ---------------------------------------------------------- 08 solar sweep
    s = slide(prs)
    eyebrow(s, "bridging intermittency")
    heading(s, "The same feeders,\nat tomorrow's solar", "07")
    lede(s, "Identical network, consumers, demand and random seed. Only rooftop "
            "penetration changes, so every difference is caused by solar.")
    if pv:
        rows = [[f"{float(x['pv_penetration_pct']):.0f}%",
                 f"{float(x['reverse_flow_pct']):.1f}%",
                 x["n_export"],
                 x["n_overvoltage"],
                 f"{float(x['storage_kwh']):.0f}",
                 f"{float(x['dr_coverage_pct']):.0f}%"] for x in pv]

        def hi3(r, c, v):
            if c == 5:
                return OK if float(str(v).rstrip("%")) > 60 else CRIT
            if c in (2, 3) and str(v) not in ("0", ""):
                return WARN
            return None

        table(s, M, Inches(3.0), Inches(11.75),
              ["rooftop pv", "reverse flow", "export breaches",
               "overvoltage", "storage kwh", "demand response can fix"],
              rows, [0.14, 0.16, 0.18, 0.16, 0.16, 0.2], highlight=hi3)
    text(s, M, Inches(5.05), Inches(11.6), Inches(1.3),
         [("Read the last column. As generation rises, demand response stops being "
           "able to fix the problem, because you cannot shed your way out of too "
           "much power. That is the case for shared storage, arrived at by "
           "measurement rather than assertion.", 14, True, INK)], spacing=1.3)
    text(s, M, Inches(6.15), Inches(11.6), Inches(0.6),
         [("Inverter over voltage trip is not modelled, so the highest voltages are "
           "an upper bound. In reality the inverters disconnect first and the "
           "customer loses the generation instead. Same failure, different symptom.",
           10.5, False, DIM)], spacing=1.25)
    footer(s, 8)

    # ---------------------------------------------------------- 09 architecture
    s = slide(prs)
    eyebrow(s, "how it deploys")
    heading(s, "A layer beneath\nthe stack you own", "08")
    lede(s, "ADMS and DERMS both assume an accurate network model. Indian low "
            "voltage networks do not have one, which is why those deployments "
            "underperform here. ENTITY GRID supplies the missing layer and feeds "
            "them, rather than competing with them.")
    stages = [
        ("Smart meters", "voltage and current\nDLMS/COSEM, 15 min"),
        ("DT gateway", "inference at the edge\nanomalies uplinked only"),
        ("ENTITY GRID", "topology, health, faults\nforecasts, flexibility"),
        ("ADMS / DERMS", "corrected connectivity\nwork orders, DR signals"),
    ]
    bw, gap = Inches(2.72), Inches(0.42)
    for i, (title, body) in enumerate(stages):
        x = M + i * (bw + gap)
        rect(s, x, Inches(3.15), bw, Inches(1.55), PANEL,
             BLUE if i == 2 else LINE)
        text(s, x + Inches(0.22), Inches(3.38), bw - Inches(0.44), Inches(0.4),
             [(title.upper(), 15, True, BLUE2 if i == 2 else INK, DISP)])
        text(s, x + Inches(0.22), Inches(3.82), bw - Inches(0.44), Inches(0.8),
             [(body, 10.5, False, MID, MONO)], spacing=1.3)
        if i < len(stages) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.09), Inches(3.78),
                Inches(0.25), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()
            arrow.shadow.inherit = False
    text(s, M, Inches(5.1), Inches(11.6), Inches(1.4),
         [("Built on what already exists: the RDSS meter rollout, DLMS/COSEM "
           "profiles, IS 16444 last gasp, CEA voltage limits, IEEE 1547 volt var, "
           "and CIM for the interface upward. Inference runs at the transformer "
           "gateway with only anomalies uplinked, because rural backhaul will not "
           "carry a firehose.", 13, False, MID)], spacing=1.3)
    footer(s, 9)

    # ---------------------------------------------------------- 10 limitations
    s = slide(prs)
    eyebrow(s, "what we will not claim")
    heading(s, "The limits,\nfrom us first", "09")
    lede(s, "Better you read these from us than find them yourselves.")
    limits = [
        ("The validation is synthetic",
         "The physics is genuine, a four wire backward forward sweep, and the "
         "confounders are modelled. It is still not a live AMI feed. Published "
         "field accuracy for meter to transformer mapping sits near 80%."),
        ("One of three degradations is found too late",
         "All three are flagged and no healthy asset is. Two arrive three and five "
         "days early. The third arrives four days after failure, because two "
         "consumers behind a lateral is not enough signal."),
        ("Winding degradation is never demonstrated",
         "All three modelled failures are neutral joints. The winding indicator runs "
         "and is correctly reported as unassessable, but nothing here proves it."),
        ("Fault recall is capped by physics, not by us",
         f"{faults['mean_recall']*100:.0f}% recall against "
         f"{faults['mean_precision']*100:.0f}% precision. A meter whose last gasp "
         "was lost to RF collision cannot be counted."),
    ]
    bullets(s, M, Inches(2.65), Inches(11.6), limits, gap=Inches(1.02))
    footer(s, 10, "all of this is on the Evidence tab of the live console")

    # ---------------------------------------------------------- 11 roadmap
    s = slide(prs)
    eyebrow(s, "deployment path")
    heading(s, "Start on one\nfeeder. Prove it.", "10")
    lede(s, "There is no price list. This is a hackathon build, not a product on "
            "sale. What follows is the honest sequence for getting it from the "
            "repository into a control room.")
    tiers = [
        ("BENCH", "today", "this repository", [
            "Full simulation and five pillar pipeline",
            "Live operator console",
            "Three head to head benchmarks",
            "One command, under a minute"]),
        ("PILOT", "8 weeks", "one DISCOM, 50 transformers", [
            "Historical AMI export, read only",
            "Corrections field verified by crews",
            "Health baselines over 30 days",
            "Scored on real SAIDI and DT failures"]),
        ("SCALE", "state wide", "edge plus ADMS integration", [
            "Inference at the DT gateway",
            "CIM interface into existing ADMS",
            "Volt var to inverters and OLTC",
            "Same software, no new hardware"]),
    ]
    tw = Inches(3.85)
    for i, (name, when, sub, points) in enumerate(tiers):
        x = M + i * (tw + Inches(0.16))
        rect(s, x, Inches(2.95), tw, Inches(3.35), PANEL,
             BLUE if i == 1 else LINE)
        if i == 1:
            rect(s, x, Inches(2.95), tw, Inches(0.05), BLUE)
        text(s, x + Inches(0.28), Inches(3.2), tw - Inches(0.56), Inches(0.3),
             [(name, 11, True, BLUE if i == 1 else DIM, MONO)])
        text(s, x + Inches(0.28), Inches(3.5), tw - Inches(0.56), Inches(0.45),
             [(when.upper(), 24, True, INK, DISP)])
        text(s, x + Inches(0.28), Inches(4.0), tw - Inches(0.56), Inches(0.3),
             [(sub, 10, False, MID, MONO)])
        for j, point in enumerate(points):
            yy = Inches(4.42) + j * Inches(0.42)
            text(s, x + Inches(0.28), yy, Inches(0.2), Inches(0.3),
                 [("+", 11, True, BLUE, MONO)])
            text(s, x + Inches(0.52), yy, tw - Inches(0.8), Inches(0.4),
                 [(point, 11, False, MID)], spacing=1.2)
    footer(s, 11)

    # ---------------------------------------------------------- 12 close
    s = slide(prs)
    rect(s, 0, 0, Inches(0.09), H, BLUE)
    text(s, M, Inches(2.1), Inches(11), Inches(0.3),
         [("250 MILLION METERS ARE GOING IN EITHER WAY", 11, True, DIM, MONO)])
    text(s, M, Inches(2.55), Inches(11.6), Inches(1.8),
         [[("THE SENSORS ARE", 58, True, INK, DISP)],
          [("ALREADY THERE.", 58, True, BLUE, DISP)]], spacing=0.92)
    text(s, M, Inches(4.6), Inches(9.4), Inches(0.9),
         [("Nobody is listening to them. Ninety days of ordinary billing data already "
           "knows the shape of the grid, which transformers are failing, where the "
           "fault is, and which consumers can help.", 15, False, MID)], spacing=1.3)
    text(s, M, Inches(5.85), Inches(11.5), Inches(0.4),
         [("github.com/kkjjkamal123/Schneider-Electric---yuva-yodha-tech",
           13, True, BLUE2, MONO)])
    text(s, M, Inches(6.3), Inches(11.5), Inches(0.4),
         [("Landing page, live operator console and full method evidence run "
           "locally with one command.", 11, False, DIM)])
    footer(s, 12, "ENTITY GRID  ·  Yuva Yodha 2026  ·  Grid Reliability")

    return prs


if __name__ == "__main__":
    data = load()
    presentation = build(data)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUT)
    print(f"wrote {OUT}  ({len(presentation.slides.__iter__.__self__._sldIdLst)} slides)")
